"""Module for reading and extracting content from presentation files."""

from pathlib import Path
from pptx import Presentation
from typing import List, Dict, Any


class PresentationReader:
    """Handles reading PowerPoint presentations and extracting text content."""
    
    def __init__(self, file_path: str):
        """
        Initialize the presentation reader.
        
        Args:
            file_path: Path to the presentation file (.pptx)
        
        Raises:
            FileNotFoundError: If the file doesn't exist
            ValueError: If the file is not a valid PowerPoint presentation
        """
        self.file_path = Path(file_path)
        
        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if self.file_path.suffix.lower() != ".pptx":
            raise ValueError(f"File must be a .pptx file, got: {self.file_path.suffix}")
        
        self.presentation = None
        self._load_presentation()
    
    def _load_presentation(self) -> None:
        """Load the PowerPoint presentation."""
        try:
            self.presentation = Presentation(self.file_path)
        except Exception as e:
            raise ValueError(f"Failed to load presentation: {str(e)}")
    
    def get_slides_content(self) -> List[Dict[str, Any]]:
        """
        Extract content from all slides.
        
        Returns:
            List of dictionaries containing slide content
        """
        slides_content = []
        
        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            slide_data = {
                "slide_number": slide_idx,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if it's a title (typically first text shape)
                    if not slide_data["title"] and hasattr(shape, "name") and "Title" in shape.name:
                        slide_data["title"] = shape.text
                    else:
                        slide_data["content"].append(shape.text)
            
            # Extract notes if available
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame.text.strip():
                    slide_data["notes"] = notes_frame.text
            
            slides_content.append(slide_data)
        
        return slides_content
    
    def get_presentation_summary(self) -> Dict[str, Any]:
        """
        Get a summary of the entire presentation.
        
        Returns:
            Dictionary with presentation metadata and content
        """
        slides = self.get_slides_content()
        
        return {
            "file_name": self.file_path.name,
            "total_slides": len(self.presentation.slides),
            "slides": slides
        }
    
    def extract_full_text(self) -> str:
        """
        Extract all text content as a single string.
        
        Returns:
            Combined text from all slides
        """
        slides = self.get_slides_content()
        full_text = []
        
        for slide in slides:
            if slide["title"]:
                full_text.append(f"Slide {slide['slide_number']}: {slide['title']}")
            
            for content in slide["content"]:
                full_text.append(content)
            
            if slide["notes"]:
                full_text.append(f"Notes: {slide['notes']}")
            
            full_text.append("")  # Empty line between slides
        
        return "\n".join(full_text)
