"""Module for generating PowerPoint slides with the executive summary."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Optional


class SlideGenerator:
    """Handles creation of new presentation slides."""
    
    def __init__(self):
        """Initialize the slide generator."""
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(7.5)
    
    def add_summary_slide(
        self,
        title: str,
        summary: str,
        subtitle: Optional[str] = None
    ) -> None:
        """
        Add an executive summary slide to the presentation.
        
        Args:
            title: The slide title
            summary: The summary content
            subtitle: Optional subtitle
        """
        # Use blank slide layout
        blank_slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(blank_slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = Inches(9)
        title_height = Inches(1)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        # Format title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add subtitle if provided
        if subtitle:
            subtitle_top = Inches(1.6)
            subtitle_height = Inches(0.5)
            
            subtitle_box = slide.shapes.add_textbox(
                title_left, subtitle_top, title_width, subtitle_height
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(18)
            subtitle_paragraph.font.color.rgb = RGBColor(100, 100, 100)
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            
            content_top = Inches(2.3)
        else:
            content_top = Inches(1.8)
        
        # Add summary content
        content_height = Inches(5)
        content_box = slide.shapes.add_textbox(
            Inches(0.8), content_top, Inches(8.4), content_height
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = summary
        
        # Format summary text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
            paragraph.level = 0
    
    def save(self, output_path: str) -> None:
        """
        Save the presentation to a file.
        
        Args:
            output_path: Path where the presentation will be saved
        
        Raises:
            ValueError: If output path is invalid
        """
        output_path = Path(output_path)
        
        if output_path.suffix.lower() != ".pptx":
            raise ValueError(f"Output file must have .pptx extension, got: {output_path.suffix}")
        
        # Create parent directories if needed
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        try:
            self.presentation.save(output_path)
        except Exception as e:
            raise Exception(f"Failed to save presentation: {str(e)}")


def create_summary_presentation(
    title: str,
    summary: str,
    output_path: str,
    subtitle: Optional[str] = None
) -> None:
    """
    Create a new presentation with an executive summary slide.
    
    Args:
        title: The slide title
        summary: The summary content
        output_path: Where to save the new presentation
        subtitle: Optional subtitle
    """
    generator = SlideGenerator()
    generator.add_summary_slide(title, summary, subtitle)
    generator.save(output_path)
